import os
import hashlib
import json
from typing import List, Optional
from fastapi import UploadFile, HTTPException
import boto3
from sqlalchemy.orm import Session
from app.models.document import Document, DocumentChunk, Embedding
from app.core.config import settings
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import openai
from sentence_transformers import SentenceTransformer

# Initialize S3 client
s3_client = boto3.client(
    's3',
    endpoint_url=settings.s3_endpoint,
    aws_access_key_id=settings.s3_access_key,
    aws_secret_access_key=settings.s3_secret_key,
    region_name=settings.s3_region
)

# Initialize embedding model
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

# Initialize OpenAI client
if settings.openai_api_key:
    openai.api_key = settings.openai_api_key


class DocumentService:
    
    @staticmethod
    async def upload_document(
        file: UploadFile,
        user_id: int,
        title: str,
        course_id: Optional[int],
        db: Session
    ) -> Document:
        """Upload and process a document"""
        
        # Validate file type
        file_extension = os.path.splitext(file.filename)[1].lower()
        if file_extension not in settings.allowed_file_types:
            raise HTTPException(
                status_code=400,
                detail=f"File type {file_extension} not allowed"
            )
        
        # Validate file size
        if file.size > settings.max_file_size:
            raise HTTPException(
                status_code=400,
                detail="File too large"
            )
        
        # Generate S3 key
        file_hash = hashlib.md5(file.filename.encode()).hexdigest()
        s3_key = f"documents/{user_id}/{file_hash}_{file.filename}"
        
        # Upload to S3
        try:
            file_content = await file.read()
            s3_client.put_object(
                Bucket=settings.s3_bucket,
                Key=s3_key,
                Body=file_content,
                ContentType=file.content_type
            )
        except Exception as e:
            raise HTTPException(
                status_code=500,
                detail=f"Failed to upload file: {str(e)}"
            )
        
        # Create document record
        document = Document(
            user_id=user_id,
            course_id=course_id,
            title=title,
            filename=file.filename,
            s3_uri=s3_key,
            mime_type=file.content_type,
            file_size=file.size,
            status="uploaded"
        )
        
        db.add(document)
        db.commit()
        db.refresh(document)
        
        # Process document asynchronously
        await DocumentService.process_document(document.id, db)
        
        return document
    
    @staticmethod
    async def process_document(document_id: int, db: Session):
        """Process document: parse, chunk, embed, generate flashcards"""
        
        document = db.query(Document).filter(Document.id == document_id).first()
        if not document:
            return
        
        try:
            # Update status to processing
            document.status = "processing"
            db.commit()
            
            # Download from S3
            response = s3_client.get_object(Bucket=settings.s3_bucket, Key=document.s3_uri)
            file_content = response['Body'].read()
            
            # Parse document
            chunks = DocumentService.parse_document(file_content, document.mime_type)
            
            # Save chunks
            for i, chunk in enumerate(chunks):
                chunk_hash = hashlib.md5(chunk['text'].encode()).hexdigest()
                
                db_chunk = DocumentChunk(
                    document_id=document.id,
                    chunk_index=i,
                    text=chunk['text'],
                    token_count=len(chunk['text'].split()),
                    chunk_hash=chunk_hash,
                    section_title=chunk.get('section_title'),
                    page_number=chunk.get('page_number')
                )
                db.add(db_chunk)
                db.flush()
                
                # Generate embedding
                embedding = DocumentService.generate_embedding(db_chunk.id, chunk['text'])
                db.add(embedding)
            
            # Update document status
            document.status = "completed"
            document.token_count = sum(len(chunk['text'].split()) for chunk in chunks)
            db.commit()
            
        except Exception as e:
            document.status = "failed"
            document.processing_error = str(e)
            db.commit()
            raise
    
    @staticmethod
    def parse_document(file_content: bytes, mime_type: str) -> List[dict]:
        """Parse document content into chunks"""
        chunks = []
        
        if mime_type == "application/pdf":
            chunks = DocumentService.parse_pdf(file_content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            chunks = DocumentService.parse_docx(file_content)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            chunks = DocumentService.parse_pptx(file_content)
        elif mime_type.startswith("image/"):
            chunks = DocumentService.parse_image(file_content)
        else:
            # Treat as text
            chunks = DocumentService.parse_text(file_content.decode('utf-8'))
        
        return chunks
    
    @staticmethod
    def parse_pdf(file_content: bytes) -> List[dict]:
        """Parse PDF content"""
        chunks = []
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        
        for page_num, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            if text.strip():
                # Split into smaller chunks if needed
                page_chunks = DocumentService.split_text_into_chunks(text, 1000)
                for i, chunk_text in enumerate(page_chunks):
                    chunks.append({
                        'text': chunk_text,
                        'page_number': page_num + 1,
                        'section_title': f"Page {page_num + 1}"
                    })
        
        return chunks
    
    @staticmethod
    def parse_docx(file_content: bytes) -> List[dict]:
        """Parse DOCX content"""
        chunks = []
        doc = DocxDocument(io.BytesIO(file_content))
        
        current_section = ""
        current_text = ""
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            
            # Check if this is a heading
            if paragraph.style.name.startswith('Heading'):
                # Save previous section
                if current_text:
                    chunks.append({
                        'text': current_text,
                        'section_title': current_section or "Introduction"
                    })
                
                current_section = text
                current_text = text + "\n"
            else:
                current_text += text + "\n"
        
        # Add final section
        if current_text:
            chunks.append({
                'text': current_text,
                'section_title': current_section or "Content"
            })
        
        return chunks
    
    @staticmethod
    def parse_pptx(file_content: bytes) -> List[dict]:
        """Parse PPTX content"""
        chunks = []
        prs = Presentation(io.BytesIO(file_content))
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                chunks.append({
                    'text': slide_text,
                    'page_number': slide_num + 1,
                    'section_title': f"Slide {slide_num + 1}"
                })
        
        return chunks
    
    @staticmethod
    def parse_image(file_content: bytes) -> List[dict]:
        """Parse image content using OCR"""
        try:
            image = Image.open(io.BytesIO(file_content))
            text = pytesseract.image_to_string(image)
            
            if text.strip():
                return [{
                    'text': text,
                    'section_title': "Image Content"
                }]
        except Exception as e:
            print(f"OCR failed: {e}")
        
        return []
    
    @staticmethod
    def parse_text(text: str) -> List[dict]:
        """Parse plain text"""
        chunks = DocumentService.split_text_into_chunks(text, 1000)
        return [{'text': chunk, 'section_title': "Text Content"} for chunk in chunks]
    
    @staticmethod
    def split_text_into_chunks(text: str, max_tokens: int) -> List[str]:
        """Split text into chunks of approximately max_tokens"""
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0
        
        for word in words:
            if current_length + len(word) + 1 > max_tokens:
                if current_chunk:
                    chunks.append(' '.join(current_chunk))
                    current_chunk = [word]
                    current_length = len(word)
                else:
                    # Single word is too long, add it anyway
                    chunks.append(word)
            else:
                current_chunk.append(word)
                current_length += len(word) + 1
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks
    
    @staticmethod
    def generate_embedding(chunk_id: int, text: str) -> Embedding:
        """Generate embedding for text chunk"""
        # Use local embedding model for now
        embedding_vector = embedding_model.encode(text)
        
        # Store in ChromaDB (simplified - in production you'd use ChromaDB client)
        vector_ref = f"chunk_{chunk_id}_{hashlib.md5(text.encode()).hexdigest()}"
        
        return Embedding(
            chunk_id=chunk_id,
            provider="local",
            model_name="all-MiniLM-L6-v2",
            vector_ref=vector_ref
        )